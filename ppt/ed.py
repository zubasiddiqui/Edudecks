

import google.generativeai as genai
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import os
import random
import re
import string
from deep_translator import GoogleTranslator

# === API KEYS ===
GENAI_API_KEY = "AIzaSyDtX3Jdl8Kz_mDyrkdHBTWB3MNNwaXFEPo"
UNSPLASH_ACCESS_KEY = "x7P3BqymG6BFDaFhnAiI_1ROJXsYt3U8xfpoLl5fEuM"

genai.configure(api_key=GENAI_API_KEY)

COLOR_PALETTE = [
    (46, 204, 113), (52, 152, 219), (155, 89, 182), (241, 196, 15),
    (230, 126, 34), (231, 76, 60), (149, 165, 166), (52, 73, 94)
]

FONT_MAP = {
    "urdu": "Jameel Noori Nastaleeq",
    "marathi": "Mangal",
    "hindi": "Mangal",
    "english": "Calibri"
}

FONT_LIST = ["Georgia", "Garamond", "Trebuchet MS", "Segoe UI", "Calibri"]


def is_dark_color(rgb_color):
    r, g, b = rgb_color
    brightness = (r * 299 + g * 587 + b * 114) / 1000
    return brightness < 128


def sanitize_filename(filename):
    valid_chars = f"-_.() {string.ascii_letters}{string.digits}"
    return ''.join(c for c in filename if c in valid_chars)


def translate_to_english(text, lang='auto'):
    try:
        return GoogleTranslator(source=lang, target='en').translate(text)
    except Exception as e:
        print(f"❌ Translation error: {e}")
        return text


def fetch_unsplash_image(subject, query, language='English'):
    search_term = f"{subject} {query}"
    if language.lower() != 'english':
        search_term = translate_to_english(search_term, lang=language.lower())
    search_term = search_term.strip().replace(" ", "+")

    url = f"https://api.unsplash.com/photos/random?query={search_term}&orientation=landscape&client_id={UNSPLASH_ACCESS_KEY}"

    try:
        response = requests.get(url)
        if response.status_code == 200:
            data = response.json()
            image_url = data['urls']['regular']
            img_data = requests.get(image_url).content
            return BytesIO(img_data)
        else:
            print(
                f"⚠️ Unsplash error: {response.status_code} - {response.text}")
    except Exception as e:
        print(f"❌ Image fetch error: {e}")
    return None


def split_into_bullets(text):
    lines = text.split('\n')
    bullets = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("•") or line.startswith("-"):
            clean = re.sub(r"^[-•\s]+", "", line)
            bullets.append(f"• {clean}")
        else:
            parts = re.split(r'(?<=[.!?])\s+', line)
            for part in parts:
                part = part.strip()
                if part:
                    bullets.append(f"• {part}")
    return bullets


def clean_title(title):
    return re.sub(r'\*\*|<u>|</u>', '', title).strip()


def generate_ppt_content(class_level, subject, topic, language='English', num_slides=5):
    model = genai.GenerativeModel('gemini-2.0-flash')
    prompt = f"""Create a {num_slides}-slide PowerPoint presentation for a {class_level}th grade {subject} class on the topic \"{topic}\".
Use only formal, culturally accurate, and age-appropriate {language} language.
Do not mix English with {language}. No Roman script or SMS-style writing.
Each slide should include:
- A meaningful, localized title (bold and underlined using **<u>Title</u>** format).
- 5-6 clear bullet points using (• or -), with examples or facts.
- Do not write paragraphs, just concise bullet points.
- Avoid any unrelated historical figures or general knowledge.
Respond exactly in this format:

SLIDE 1: [Title]
[• Bullet or - Bullet or sentence]

SLIDE 2: [Title]
[• Bullet]
[...]"""
    try:
        response = model.generate_content(prompt)
        return format_ppt_output(response.text)
    except Exception as e:
        return f"❌ Error generating content: {str(e)}"


def format_ppt_output(raw_text):
    slides = []
    current_slide = {}
    for line in raw_text.split('\n'):
        if line.strip().startswith('SLIDE'):
            if current_slide:
                slides.append(current_slide)
            parts = line.split(':', 1)
            current_slide = {'title': clean_title(
                parts[1].strip()), 'content': []}
        elif line.strip():
            current_slide['content'].append(line.strip())
    if current_slide:
        slides.append(current_slide)
    return slides


def create_powerpoint(slides, topic, class_level, subject, language="English"):
    prs = Presentation()
    language_key = language.lower().strip()
    rtl = True if language_key == 'urdu' else False

    bg_rgb = random.choice(COLOR_PALETTE)
    font_name = FONT_MAP.get(language_key, random.choice(FONT_LIST))
    is_dark = is_dark_color(bg_rgb)
    font_color = RGBColor(255, 255, 255) if is_dark else RGBColor(0, 0, 0)

    # === Title Slide ===
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(*bg_rgb)

    title_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(1.5), prs.slide_width - Inches(2), Inches(2))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = title_tf.paragraphs[0]
    p.text = slides[0]['title']
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    title_len = len(slides[0]['title'])
    if title_len > 100:
        p.font.size = Pt(24)
    elif title_len > 60:
        p.font.size = Pt(28)
    else:
        p.font.size = Pt(36)

    if slides[0]['content']:
        subtitle = slides[0]['content'][0]
        subtitle_box = title_slide.shapes.add_textbox(
            Inches(1), Inches(3.5), prs.slide_width - Inches(2), Inches(2))
        sub_tf = subtitle_box.text_frame
        sub_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        sub_tf.word_wrap = True
        sub_tf.text = subtitle[:250]
        sp = sub_tf.paragraphs[0]
        sp.font.size = Pt(22)
        sp.font.color.rgb = font_color
        sp.font.name = font_name
        sp.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    # === Content Slides ===
    for slide in slides[1:]:
        ppt_slide = prs.slides.add_slide(prs.slide_layouts[6])
        ppt_slide.background.fill.solid()
        ppt_slide.background.fill.fore_color.rgb = RGBColor(*bg_rgb)

        slide_width = prs.slide_width
        image_width = slide_width * 0.35
        text_width = slide_width * 0.60
        margin = Inches(0.4)

        title_box = ppt_slide.shapes.add_textbox(
            margin, Inches(0.3), text_width - margin, Inches(1))
        title_tf = title_box.text_frame
        title_tf.clear()
        p = title_tf.paragraphs[0]
        p.text = slide['title']
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.underline = True
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

        content_box = ppt_slide.shapes.add_textbox(margin, Inches(
            1.2), text_width - margin, prs.slide_height - Inches(1.7))
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        content_tf.margin_top = Inches(0.1)
        content_tf.margin_bottom = Inches(0.1)
        content_tf.margin_left = Inches(0.1)
        content_tf.margin_right = Inches(0.1)

        points = []
        for raw_line in slide['content']:
            points.extend(split_into_bullets(raw_line))

        point_count = len(points)
        if point_count <= 4:
            font_size = Pt(26)
        elif point_count <= 6:
            font_size = Pt(22)
        else:
            font_size = Pt(18)

        for bullet in points[:6]:
            para = content_tf.add_paragraph()
            para.text = bullet
            para.font.size = font_size
            para.font.color.rgb = font_color
            para.font.name = font_name
            para.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
            para.line_spacing = Pt(font_size.pt + 6)  # Line spacing added here

        img_stream = fetch_unsplash_image(subject, slide['title'], language)
        if img_stream:
            try:
                img_path = f"temp_img_{slide['title'].replace(' ', '_')}.jpg"
                with open(img_path, "wb") as f:
                    f.write(img_stream.read())
                ppt_slide.shapes.add_picture(
                    img_path,
                    left=slide_width - image_width,
                    top=Inches(1.2),
                    width=image_width - Inches(0.4),
                    height=prs.slide_height - Inches(1.7)
                )
                os.remove(img_path)
            except Exception as e:
                print(f"⚠️ Could not add image: {e}")

        footer = ppt_slide.shapes.add_textbox(
            margin, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5))
        footer_tf = footer.text_frame
        footer_text = {
            "english": "Generated by AI Slide Generator",
            "hindi": "एआई स्लाइड जेनरेटर द्वारा निर्मित",
            "marathi": "एआय स्लाइड जनरेटर द्वारे तयार केले",
            "urdu": "اے آئی سلائیڈ جنریٹر کے ذریعہ تیار کردہ"
        }.get(language_key, "Generated by AI Slide Generator")
        footer_tf.text = footer_text
        footer_tf.paragraphs[0].font.size = Pt(12)
        footer_tf.paragraphs[0].font.color.rgb = font_color
        footer_tf.paragraphs[0].font.name = font_name
        footer_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    filename = sanitize_filename(
        f"Class{class_level}_{subject}_{topic}_presentation.pptx".replace(" ", "_"))
    prs.save(filename)
    try:
        print(f"✅ PowerPoint saved as: {filename}")
    except UnicodeEncodeError:
        print(f"PowerPoint saved as: {filename}")


def run_generator():
    try:
        print("\n🎓 AI Slide Generator")
    except UnicodeEncodeError:
        print("\nAI Slide Generator")
    class_level = input("Grade level (e.g. 5): ")
    subject = input("Subject: ")
    topic = input("Topic: ")
    language = input("Language [default English]: ") or "English"
    try:
        num_slides = int(input("Number of slides [3-10, default 5]: ") or 5)
        if num_slides < 3 or num_slides > 10:
            raise ValueError
    except ValueError:
        print("❌ Number of slides must be between 3 and 10.")
        return

    content = generate_ppt_content(
        class_level, subject, topic, language, num_slides)
    if isinstance(content, list):
        create_powerpoint(content, topic, class_level, subject, language)
    else:
        print(content)


if __name__ == "__main__":
    run_generator()
